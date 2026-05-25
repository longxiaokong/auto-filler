"""文件处理模块 - 增强版"""
import hashlib
import importlib.util
import logging
import os
import re
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

import markdown_it
from openpyxl import load_workbook
from pypdf import PdfReader
from pypdf.errors import PdfReadError

from app.config import FileOCRConfig, settings


def _detect_inference_device(device_config: str = "auto") -> str:
    """自动检测推理设备，避免为拿设备信息导入完整 retrieval 模块。"""
    if device_config != "auto":
        return device_config

    try:
        import torch

        if torch.cuda.is_available():
            return "cuda"
        if hasattr(torch.backends, "mps") and torch.backends.mps.is_available():
            return "mps"
        return "cpu"
    except ImportError:
        logger.warning("PyTorch not available, using CPU")
        return "cpu"


class OCRProcessor:
    """OCR 处理器 - 基于 RapidOCR 的轻量本地实现"""

    def __init__(self):
        self.config = getattr(settings.file, "ocr", None) or FileOCRConfig()
        if not hasattr(settings.file, "ocr"):
            logger.warning("settings.file 未配置 ocr，使用默认 OCR 配置")
        self.enabled = self.config.enabled
        self._engine = None
        self._load_attempted = False
        self._load_error: Optional[str] = None

    def _ensure_engine(self):
        """惰性加载 OCR 引擎，避免进程启动时阻塞模型初始化。"""
        if not self.enabled:
            return None
        if self._engine is not None:
            return self._engine
        if self._load_attempted:
            return None

        self._load_attempted = True

        if self.config.provider.lower() != "rapidocr":
            self._load_error = f"不支持的 OCR provider: {self.config.provider}"
            logger.warning(self._load_error)
            return None

        if importlib.util.find_spec("rapidocr") is None:
            self._load_error = "rapidocr 未安装，图片 OCR 已自动降级"
            logger.warning(self._load_error)
            return None

        try:
            from rapidocr import (
                EngineType,
                LangCls,
                LangDet,
                LangRec,
                ModelType,
                OCRVersion,
                RapidOCR,
            )

            # 仅覆盖当前项目需要的关键参数，其余保持 RapidOCR 默认值。
            params = {
                "Global.text_score": self.config.text_score,
                "Global.use_cls": self.config.use_cls,
                "Global.log_level": self.config.log_level,
                "Det.engine_type": EngineType(self.config.engine_type),
                "Det.lang_type": LangDet(self.config.det_lang_type),
                "Det.model_type": ModelType(self.config.model_type),
                "Det.ocr_version": OCRVersion(self.config.ocr_version),
                "Cls.engine_type": EngineType(self.config.engine_type),
                "Cls.lang_type": LangCls("ch"),
                "Cls.model_type": ModelType(self.config.model_type),
                "Rec.engine_type": EngineType(self.config.engine_type),
                "Rec.lang_type": LangRec(self.config.rec_lang_type),
                "Rec.model_type": ModelType(self.config.model_type),
                "Rec.ocr_version": OCRVersion(self.config.ocr_version),
            }
            if self.config.use_coreml:
                coreml_cache_dir = str(Path(self.config.coreml_cache_dir).resolve())
                Path(coreml_cache_dir).mkdir(parents=True, exist_ok=True)
                params["EngineConfig.onnxruntime.use_coreml"] = True
                params["EngineConfig.onnxruntime.coreml_ep_cfg.ModelCacheDirectory"] = coreml_cache_dir
            self._engine = RapidOCR(params=params)
            logger.info(
                "OCR 引擎初始化完成: provider=%s engine=%s det=%s/%s rec=%s/%s",
                self.config.provider,
                self.config.engine_type,
                self.config.ocr_version,
                self.config.model_type,
                self.config.ocr_version,
                self.config.model_type,
            )
            return self._engine
        except Exception as e:
            self._load_error = str(e)
            logger.warning("OCR 引擎初始化失败，图片 OCR 已自动降级: %s", e)
            return None

    def _extract_text(self, result: Any) -> str:
        """兼容 RapidOCR 新旧结果结构，提取纯文本。"""
        if result is None:
            return ""

        lines: List[str] = []
        threshold = float(self.config.text_score)

        txts = getattr(result, "txts", None)
        scores = getattr(result, "scores", None)
        if txts:
            for idx, text in enumerate(txts):
                if not isinstance(text, str):
                    text = str(text)
                score = None
                if scores and idx < len(scores):
                    score = scores[idx]
                if text.strip() and (score is None or float(score) >= threshold):
                    lines.append(text.strip())
            return "\n".join(lines)

        if isinstance(result, (list, tuple)):
            for item in result:
                if not isinstance(item, (list, tuple)) or len(item) < 2:
                    continue
                text = item[1]
                score = item[2] if len(item) > 2 else None
                if text and (score is None or float(score) >= threshold):
                    lines.append(str(text).strip())
            return "\n".join(line for line in lines if line)

        return ""

    def process_image(self, image_path: str) -> str:
        """处理图片，提取文字"""
        engine = self._ensure_engine()
        if engine is None:
            return ""

        try:
            result = engine(image_path)
            return self._extract_text(result)
        except Exception as e:
            logger.warning("OCR 识别失败，文件=%s: %s", image_path, e)
            return ""

    def process_image_bytes(self, image_bytes: bytes) -> str:
        """处理图片字节，提取文字"""
        if not self.enabled:
            return ""

        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(image_bytes)
                tmp_path = tmp.name
            return self.process_image(tmp_path)
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.unlink(tmp_path)


class FileProcessor:
    """文件处理器 - 支持多种文档格式"""

    def __init__(self):
        self.chunk_size = settings.file.chunk_size
        self.chunk_overlap = settings.file.chunk_overlap
        self.parent_chunk_size = settings.file.parent_chunk_size
        self.strategy = settings.file.chunk_strategy
        self.upload_dir = settings.file.upload_dir
        self.supported_formats = settings.file.supported_formats

        # OCR处理器
        self.ocr = OCRProcessor()

        # 语义分块模型缓存（避免每次调用都重新加载）
        self._semantic_model = None

        # 设备检测（用于语义分块）
        self._device = _detect_inference_device(settings.vector.device)

        # 确保上传目录存在
        Path(self.upload_dir).mkdir(parents=True, exist_ok=True)

    def process_file(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理文件并返回文本块"""
        ext = Path(file_path).suffix.lower()

        # 根据文件类型选择处理方法
        if ext == ".pdf":
            return self._process_pdf(file_path, project_id)
        elif ext in [".docx", ".doc"]:
            return self._process_docx(file_path, project_id)
        elif ext == ".md":
            return self._process_markdown(file_path, project_id)
        elif ext == ".txt":
            return self._process_txt(file_path, project_id)
        elif ext in [".xlsx", ".xls"]:
            return self._process_excel(file_path, project_id)
        elif ext in [".pptx", ".ppt"]:
            return self._process_pptx(file_path, project_id)
        elif ext == ".csv":
            return self._process_csv(file_path, project_id)
        elif ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"]:
            return self._process_image(file_path, project_id)
        elif ext == ".html":
            return self._process_html(file_path, project_id)
        elif ext == ".json":
            return self._process_json(file_path, project_id)
        else:
            raise ValueError(f"不支持的文件类型: {ext}")

    def _process_pdf(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理 PDF 文件 - 增强版：提取文本和表格"""
        try:
            reader = PdfReader(file_path)
            text = ""
            images_count = 0
            tables_count = 0
            pages_text = []

            for i, page in enumerate(reader.pages):
                # 提取文本
                page_text = page.extract_text()
                if page_text:
                    text += f"\n## 第 {i+1} 页\n"
                    text += page_text + "\n"
                    pages_text.append(page_text)

                # 检查是否有图片（用于OCR预留）- 修复bug，添加异常处理
                try:
                    if "/Resources" in page and "/XObject" in page["/Resources"]:
                        xobjects = page["/Resources"]["/XObject"].get_object()
                        for obj in xobjects:
                            if xobjects[obj]["/Subtype"] == "/Image":
                                images_count += 1
                except Exception:
                    pass  # 忽略资源访问错误

            # TODO: 表格提取 - 使用 pdfplumber 或 Camelot
            # if tables_count == 0 and images_count > 0:
            #     # 纯图PDF，走OCR流程
            #     pass

            chunks = self._chunk_text(text)

            return {
                "file_id": self._generate_file_id(file_path),
                "filename": Path(file_path).name,
                "file_type": "pdf",
                "total_pages": len(reader.pages),
                "total_chars": len(text),
                "images_count": images_count,
                "tables_count": tables_count,
                "chunks": chunks,
                "project_id": project_id,
            }

        except PdfReadError as e:
            raise ValueError(f"PDF 文件读取失败: {str(e)}")

    def _process_docx(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理 Word 文件 - 支持 docx 和 doc 格式"""
        import zipfile
        from docx import Document
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        ext = Path(file_path).suffix.lower()
        text = ""
        paragraphs_count = 0
        tables_count = 0
        images_count = 0

        # .doc 格式使用多策略处理（跨平台）
        if ext == ".doc":
            import subprocess
            import platform
            import shutil
            import tempfile
            import re

            text = ""
            last_error = None

            # 策略1: 尝试 antiword (Linux/macOS)
            antiword_path = shutil.which("antiword")
            if antiword_path:
                try:
                    result = subprocess.run(
                        ["antiword", file_path],
                        capture_output=True,
                        text=True,
                        timeout=60
                    )
                    if result.returncode == 0 and result.stdout:
                        text = result.stdout.strip()
                except Exception as e:
                    last_error = e

            # 策略2: 尝试 textutil (macOS only)
            if not text and platform.system() == "Darwin":
                tmp_path = None
                try:
                    # 创建临时文件
                    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False, mode='w') as tmp:
                        tmp_path = tmp.name
                    result = subprocess.run(
                        ["textutil", "-convert", "txt", "-output", tmp_path, file_path],
                        capture_output=True,
                        text=True,
                        timeout=60
                    )
                    if result.returncode == 0:
                        with open(tmp_path, 'r', encoding='utf-8') as f:
                            text = f.read().strip()
                        os.unlink(tmp_path)
                        tmp_path = None
                except Exception as e:
                    last_error = e
                    if tmp_path and os.path.exists(tmp_path):
                        os.unlink(tmp_path)

            # 策略3: 尝试 libreoffice (跨平台)
            if not text:
                libreoffice_path = shutil.which("soffice") or shutil.which("libreoffice")
                if libreoffice_path:
                    try:
                        # 使用 libreoffice 命令行转换
                        with tempfile.TemporaryDirectory() as tmpdir:
                            result = subprocess.run(
                                [libreoffice_path, "--headless", "--convert-to", "txt", "--outdir", tmpdir, file_path],
                                capture_output=True,
                                text=True,
                                timeout=120
                            )
                            if result.returncode == 0:
                                # 读取转换后的文件
                                base_name = Path(file_path).stem
                                txt_file = Path(tmpdir) / f"{base_name}.txt"
                                if txt_file.exists():
                                    with open(txt_file, 'r', encoding='utf-8', errors='ignore') as f:
                                        text = f.read().strip()
                    except Exception as e:
                        last_error = e

            # 策略4: 尝试 pandoc (跨平台)
            if not text:
                pandoc_path = shutil.which("pandoc")
                if pandoc_path:
                    try:
                        result = subprocess.run(
                            [pandoc_path, "-t", "plain", file_path],
                            capture_output=True,
                            text=True,
                            timeout=60
                        )
                        if result.returncode == 0 and result.stdout:
                            text = result.stdout.strip()
                    except Exception as e:
                        last_error = e

            # 策略5: 尝试读取 OLE 文件内部 XML (仅作为最后手段)
            if not text:
                try:
                    import olefile
                    ole = olefile.OleFileIO(file_path)
                    if 'WordDocument' in ole.listdir():
                        # 尝试提取原始文本数据
                        stream = ole.open_stream('WordDocument')
                        data = stream.read()
                        # 简单的 Unicode 文本提取
                        text_parts = []
                        i = 0
                        while i < len(data) - 2:
                            # 尝试查找 Unicode 文本模式
                            if data[i:i+2] == b'\x00':
                                try:
                                    char = chr(data[i+1]) if data[i+1] < 128 else None
                                    if char and char.isprintable():
                                        text_parts.append(char)
                                except Exception:
                                    pass
                            i += 2
                        text = ''.join(text_parts)
                        # 清理提取的文本
                        text = re.sub(r'[^\x20-\x7E\n\r\t]', '', text)
                        text = text.strip()
                    ole.close()
                except Exception as e:
                    last_error = e

            if not text:
                raise ValueError(f"无法读取 .doc 文件: {file_path}, 请安装 antiword/libreoffice/pandoc: {last_error}")

            # 计算段落数量
            paragraphs_count = len([p for p in text.split('\n') if p.strip()])

        else:
            # .docx 格式使用 python-docx 处理
            try:
                doc = Document(file_path)
                text = ""

                # 按文档顺序提取所有内容（段落+表格）
                for element in doc.element.body:
                    if isinstance(element, CT_P):
                        # 段落
                        para = Paragraph(element, doc)
                        para_text = para.text.strip()
                        if para_text:
                            # 检查是否为标题
                            if para.style.name.startswith('Heading'):
                                text += f"\n## {para_text}\n"
                            else:
                                text += para_text + "\n"
                            paragraphs_count += 1

                    elif isinstance(element, CT_Tbl):
                        # 表格
                        table = Table(element, doc)
                        table_text = self._extract_table_text(table)
                        if table_text:
                            text += f"\n### 表格 {tables_count + 1}\n"
                            text += table_text + "\n"
                            tables_count += 1

                # 提取图片（预留）
                images_count = len(doc.part.related_parts)

            except Exception as e:
                # .docx 解析失败
                if not text:
                    raise ValueError(f"无法读取 Word 文件: {file_path}, 错误: {str(e)}")

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "docx",
            "total_paragraphs": paragraphs_count,
            "total_tables": tables_count,
            "images_count": images_count,
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _extract_table_text(self, table) -> str:
        """提取表格文本"""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ')
                cells.append(cell_text)
            if cells:
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _process_pptx(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理 PowerPoint 文件 - 支持 .pptx 和 .ppt 格式"""
        ext = Path(file_path).suffix.lower()

        # .ppt 格式使用多策略处理（类似 .doc）
        if ext == ".ppt":
            return self._process_ppt_legacy(file_path, project_id)

        try:
            from pptx import Presentation
        except ImportError:
            # 如果没有 pptx 库，尝试从 XML 提取
            try:
                import zipfile
                text = ""
                with zipfile.ZipFile(file_path, 'r') as zf:
                    # 读取所有幻灯片 XML
                    for name in sorted(zf.namelist()):
                        if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                            with zf.open(name) as xml_file:
                                import xml.etree.ElementTree as ET
                                content = xml_file.read()
                                # 提取所有文本
                                import re
                                texts = re.findall(r'<a:t>([^<]+)</a:t>', content.decode('utf-8'))
                                if texts:
                                    slide_num = name.replace('ppt/slides/slide', '').replace('.xml', '')
                                    text += f"\n## 幻灯片 {slide_num}\n"
                                    text += "\n".join(texts) + "\n"

                if not text:
                    raise ImportError("No pptx library")

                chunks = self._chunk_text(text)
                return {
                    "file_id": self._generate_file_id(file_path),
                    "filename": Path(file_path).name,
                    "file_type": "pptx",
                    "total_slides": 0,
                    "total_shapes": 0,
                    "total_tables": 0,
                    "total_chars": len(text),
                    "chunks": chunks,
                    "project_id": project_id,
                }
            except Exception as e:
                raise ValueError(f"需要安装 python-pptx 库来读取 PPTX 文件: {file_path}")

        # 使用 python-pptx 解析
        prs = Presentation(file_path)
        text = ""
        slides_count = len(prs.slides)
        shapes_count = 0
        tables_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"\n## 幻灯片 {slide_num}\n"

            for shape in slide.shapes:
                shapes_count += 1

                # 提取表格
                if shape.has_table:
                    tables_count += 1
                    table_text = self._extract_ppt_table_text(shape.table)
                    slide_text += f"\n### 表格\n{table_text}\n"

                # 提取文本框
                elif shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            # 检查是否为标题（根据级别）
                            if paragraph.level == 0:
                                slide_text += f"**{para_text}**\n"
                            else:
                                slide_text += para_text + "\n"

            text += slide_text

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "pptx",
            "total_slides": slides_count,
            "total_shapes": shapes_count,
            "total_tables": tables_count,
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _extract_ppt_table_text(self, table) -> str:
        """提取PPT表格文本"""
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ')
                cells.append(cell_text)
            if cells:
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _process_ppt_legacy(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理旧版 .ppt 文件（二进制格式）- 多策略降级"""
        import subprocess
        import shutil
        import tempfile

        text = ""
        last_error = None

        # 策略1: catppt (catdoc 套件)
        catppt_path = shutil.which("catppt")
        if catppt_path:
            try:
                result = subprocess.run(
                    ["catppt", file_path],
                    capture_output=True, text=True, timeout=60
                )
                if result.returncode == 0 and result.stdout:
                    text = result.stdout.strip()
            except Exception as e:
                last_error = e

        # 策略2: libreoffice (跨平台)
        if not text:
            libreoffice_path = shutil.which("soffice") or shutil.which("libreoffice")
            if libreoffice_path:
                try:
                    with tempfile.TemporaryDirectory() as tmpdir:
                        result = subprocess.run(
                            [libreoffice_path, "--headless", "--convert-to", "txt", "--outdir", tmpdir, file_path],
                            capture_output=True, text=True, timeout=120
                        )
                        if result.returncode == 0:
                            base_name = Path(file_path).stem
                            txt_file = Path(tmpdir) / f"{base_name}.txt"
                            if txt_file.exists():
                                with open(txt_file, 'r', encoding='utf-8', errors='ignore') as f:
                                    text = f.read().strip()
                except Exception as e:
                    last_error = e

        # 策略3: OLE SummaryInformation 提取标题/主题等元数据
        if not text:
            try:
                import olefile
                ole = olefile.OleFileIO(file_path)
                meta = ole.get_metadata()
                parts = []
                for attr in ['title', 'subject', 'author', 'comments', 'keywords']:
                    val = getattr(meta, attr, None)
                    if val:
                        if isinstance(val, bytes):
                            val = val.decode('utf-8', errors='ignore')
                        val = str(val).strip()
                        if val:
                            parts.append(f"{attr}: {val}")
                ole.close()
                if parts:
                    text = '\n'.join(parts)
            except Exception as e:
                last_error = e

        # 如果提取失败或文本太短，生成占位说明
        if not text or len(text) < 10:
            text = f"[PPT 文件 - 无法提取文本内容: {Path(file_path).name}]"

        # 文本质量检验：过滤乱码
        text = self._validate_text_quality(text)

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "ppt",
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _validate_text_quality(self, text: str) -> str:
        """检验文本质量，过滤乱码"""
        if not text:
            return text
        # 统计可识别字符比例（ASCII可打印 + CJK + 常见标点）
        total = len(text)
        if total == 0:
            return text
        readable = sum(
            1 for c in text
            if c.isascii() and c.isprintable()
            or '\u4e00' <= c <= '\u9fff'  # CJK
            or '\u3000' <= c <= '\u303f'  # CJK 标点
            or '\uff00' <= c <= '\uffef'  # 全角字符
            or c in '\n\r\t'
        )
        ratio = readable / total
        if ratio < 0.5:
            # 超过一半是乱码，丢弃并返回占位文本
            return ""
        return text

    def _process_markdown(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理 Markdown 文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        # 解析 Markdown，保留结构
        md = markdown_it.MarkdownIt()
        tokens = md.parse(content)

        # 提取纯文本，保留标题层级（h1→#, h2→##, ...）
        plain_text = ""
        pending_heading_prefix = ""
        for token in tokens:
            if token.type == "heading_open":
                level = int(token.tag[1]) if token.tag and len(token.tag) == 2 and token.tag[1:].isdigit() else 1
                pending_heading_prefix = "#" * level + " "
            elif token.type == "inline" and token.content:
                if pending_heading_prefix:
                    plain_text += pending_heading_prefix + token.content + "\n"
                    pending_heading_prefix = ""
                else:
                    plain_text += token.content + "\n"
            elif token.type == "heading_close":
                pending_heading_prefix = ""

        if not plain_text:
            plain_text = content

        chunks = self._chunk_text(plain_text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "markdown",
            "total_chars": len(content),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _process_txt(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理文本文件"""
        # 尝试多种编码
        text = ""
        for encoding in ['utf-8', 'gbk', 'gb2312', 'latin1']:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                break
            except UnicodeDecodeError:
                continue

        if not text:
            raise ValueError(f"无法解码文本文件: {file_path}")

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "txt",
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _process_excel(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理 Excel 文件 - 支持 xlsx 和 xls 格式"""
        ext = Path(file_path).suffix.lower()

        text = ""
        total_rows = 0
        total_tables = 0
        total_sheets = 0

        if ext == ".xls":
            # 使用 xlrd 处理旧版 .xls 格式
            try:
                import xlrd
                wb = xlrd.open_workbook(file_path)
                total_sheets = wb.nsheets

                for sheet_idx in range(wb.nsheets):
                    sheet = wb.sheet_by_index(sheet_idx)
                    sheet_name = sheet.name or f"Sheet{sheet_idx + 1}"
                    text += f"\n## 工作表: {sheet_name}\n"

                    rows_data = []
                    for row_idx in range(sheet.nrows):
                        row = []
                        for col_idx in range(sheet.ncols):
                            cell = sheet.cell(row_idx, col_idx)
                            if cell.ctype == xlrd.XL_CELL_EMPTY:
                                row.append("")
                            else:
                                row.append(str(cell.value))
                        rows_data.append(tuple(row))

                    if rows_data:
                        table_text = self._extract_excel_table(rows_data)
                        text += table_text + "\n"
                        total_tables += 1
                        total_rows += len(rows_data)

            except ImportError:
                raise ValueError(f"需要安装 xlrd 库来读取 .xls 文件: {file_path}")
        else:
            # 使用 openpyxl 处理 .xlsx 格式
            wb = load_workbook(file_path, data_only=True, read_only=False)
            total_sheets = len(wb.sheetnames)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"\n## 工作表: {sheet_name}\n"

                rows_data = []
                for row in sheet.iter_rows(values_only=True):
                    rows_data.append(row)

                if rows_data:
                    table_text = self._extract_excel_table(rows_data)
                    text += table_text + "\n"
                    total_tables += 1
                    total_rows += len(rows_data)

            wb.close()

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "xlsx",
            "total_sheets": total_sheets,
            "total_tables": total_tables,
            "total_rows": total_rows,
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _extract_excel_table(self, rows_data: List[Tuple]) -> str:
        """提取Excel表格数据"""
        if not rows_data:
            return ""

        lines = []
        for row in rows_data:
            cells = []
            for cell in row:
                if cell is not None:
                    cell_str = str(cell)
                    cells.append(cell_str)
                else:
                    cells.append("")
            # 只添加非空行
            if any(cells):
                lines.append(" | ".join(cells))

        return "\n".join(lines)

    def _process_csv(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理 CSV 文件"""
        import csv

        text = ""
        rows_count = 0

        # 尝试多种编码
        for encoding in ['utf-8', 'gbk', 'gb2312', 'latin1']:
            try:
                with open(file_path, 'r', encoding=encoding, newline='') as f:
                    reader = csv.reader(f)
                    for row in reader:
                        if any(cell.strip() for cell in row):
                            text += " | ".join(row) + "\n"
                            rows_count += 1
                break
            except UnicodeDecodeError:
                continue

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "csv",
            "total_rows": rows_count,
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _process_image(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理图片文件 - 使用OCR"""
        filename = Path(file_path).name
        ocr_text = ""

        # 如果OCR启用，使用OCR提取文字
        if self.ocr.enabled:
            ocr_text = self.ocr.process_image(file_path)
            if ocr_text:
                text = f"文件名: {filename}\nOCR文本:\n{ocr_text}"
            else:
                text = f"文件名: {filename}\n[图片OCR未识别到文本]"
        else:
            # OCR未启用，记录需要处理
            text = f"文件名: {filename}\n[图片文件 - 需要OCR处理]"

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "image",
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _process_html(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理 HTML 文件"""
        from html.parser import HTMLParser

        class HTMLTextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self.skip = False

            def handle_starttag(self, tag, attrs):
                if tag in ['script', 'style']:
                    self.skip = True
                elif tag in ['p', 'div', 'br', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    self.text.append('\n')

            def handle_endtag(self, tag):
                if tag in ['script', 'style']:
                    self.skip = False
                elif tag in ['p', 'div']:
                    self.text.append('\n')

            def handle_data(self, data):
                if not self.skip:
                    self.text.append(data)

        with open(file_path, "r", encoding="utf-8") as f:
            html_content = f.read()

        extractor = HTMLTextExtractor()
        extractor.feed(html_content)
        text = "".join(extractor.text)

        # 清理文本
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = text.strip()

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "html",
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _process_json(self, file_path: str, project_id: Optional[str] = None) -> Dict[str, Any]:
        """处理 JSON 文件"""
        import json

        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        # 递归提取JSON中的所有文本
        text = self._extract_json_text(data, indent=0)

        chunks = self._chunk_text(text)

        return {
            "file_id": self._generate_file_id(file_path),
            "filename": Path(file_path).name,
            "file_type": "json",
            "total_chars": len(text),
            "chunks": chunks,
            "project_id": project_id,
        }

    def _extract_json_text(self, obj, indent: int = 0) -> str:
        """递归提取JSON文本"""
        text = ""
        prefix = "  " * indent

        if isinstance(obj, dict):
            for key, value in obj.items():
                if isinstance(value, (dict, list)):
                    text += f"{prefix}{key}:\n{self._extract_json_text(value, indent + 1)}"
                else:
                    text += f"{prefix}{key}: {value}\n"
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                if isinstance(item, (dict, list)):
                    text += f"{prefix}[{i}]\n{self._extract_json_text(item, indent + 1)}"
                else:
                    text += f"{prefix}[{i}]: {item}\n"

        return text

    def _chunk_text(self, text: str) -> List[Dict[str, Any]]:
        """根据配置策略将文本分块"""
        if not text:
            return []
        text = self._clean_text(text)
        logger.debug("分块策略=%s 文本长度=%d", self.strategy, len(text))
        if self.strategy == "recursive":
            chunks = self._chunk_text_recursive(text)
        elif self.strategy == "semantic":
            chunks = self._chunk_text_semantic(text)
        elif self.strategy == "hierarchical":
            chunks = self._chunk_text_hierarchical(text)
        else:
            raise ValueError(f"Unknown chunk_strategy: {self.strategy}")
        logger.debug("分块完成，共 %d 个 chunk", len(chunks))
        return chunks

    # ── 递归分块 ─────────────────────────────────────────────────────────────
    def _chunk_text_recursive(
        self,
        text: str,
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
    ) -> List[Dict[str, Any]]:
        """递归分块：按优先级分隔符拆分，保证句子完整性"""
        if chunk_size is None:
            chunk_size = self.chunk_size
        if chunk_overlap is None:
            chunk_overlap = self.chunk_overlap

        separators = ["\n\n", "\n", "。", "！", "？", ".", "!", "?", " ", ""]

        def split_by_separators(s: str, seps: List[str]) -> List[str]:
            if not seps or len(s) <= chunk_size:
                return [s] if s.strip() else []
            sep = seps[0]
            parts = s.split(sep) if sep else list(s)
            result: List[str] = []
            current = ""
            for part in parts:
                segment = (current + sep + part) if current else part
                if len(segment) <= chunk_size:
                    current = segment
                else:
                    if current:
                        result.append(current)
                    # 子片段仍然太长时递归处理
                    if len(part) > chunk_size:
                        result.extend(split_by_separators(part, seps[1:]))
                        current = ""
                    else:
                        current = part
            if current:
                result.append(current)
            return result

        raw_chunks = split_by_separators(text, separators)

        # 合并过短的片段，添加 overlap
        merged: List[str] = []
        buf = ""
        for seg in raw_chunks:
            if not seg.strip():
                continue
            candidate = (buf + seg) if not buf else (buf + "\n" + seg)
            if len(candidate) <= chunk_size:
                buf = candidate
            else:
                if buf:
                    merged.append(buf)
                buf = seg
        if buf:
            merged.append(buf)

        # 构造带 overlap 的最终块
        result: List[Dict[str, Any]] = []
        for i, content in enumerate(merged):
            overlap_prefix = ""
            if i > 0 and chunk_overlap > 0:
                prev = merged[i - 1]
                overlap_prefix = prev[-chunk_overlap:] if len(prev) > chunk_overlap else prev
            final_content = (overlap_prefix + "\n" + content).strip() if overlap_prefix else content
            # 截断到 chunk_size，保留末尾（即实际内容），丢弃多余的 overlap 前缀
            if len(final_content) > chunk_size:
                final_content = final_content[-chunk_size:]
            result.append({
                "chunk_id": f"chunk_{i}",
                "content": final_content,
            })
        return result

    # ── 语义分块 ─────────────────────────────────────────────────────────────
    def _chunk_text_semantic(self, text: str) -> List[Dict[str, Any]]:
        """语义分块：在相邻句子嵌入相似度骤降处断块"""
        import re as _re

        # 按句子边界切分
        sentence_endings = _re.compile(r"(?<=[。！？.!?\n])\s*")
        sentences = [s.strip() for s in sentence_endings.split(text) if s.strip()]
        if len(sentences) <= 1:
            # 退化到 recursive
            return self._chunk_text_recursive(text)

        # 尝试使用 embedding 模型计算相似度（模型缓存为实例属性，避免重复加载）
        try:
            import numpy as np
            if self._semantic_model is None:
                from sentence_transformers import SentenceTransformer
                from app.config import settings as _settings
                from app.core.retrieval import resolve_local_sentence_transformer
                local_model_path = resolve_local_sentence_transformer(
                    _settings.vector.embedding_model
                )
                self._semantic_model = SentenceTransformer(
                    local_model_path,
                    device=self._device,
                )
                logger.debug(
                    "语义分块模型加载完成: %s -> %s (设备: %s)",
                    _settings.vector.embedding_model,
                    local_model_path,
                    self._device,
                )
            model = self._semantic_model
            embeddings = model.encode(
                sentences,
                device=self._device,
                show_progress_bar=False
            )
            sims = [
                float(np.dot(embeddings[i], embeddings[i + 1]) /
                      (np.linalg.norm(embeddings[i]) * np.linalg.norm(embeddings[i + 1]) + 1e-10))
                for i in range(len(embeddings) - 1)
            ]
            mean_sim = float(np.mean(sims))
            std_sim = float(np.std(sims))
            threshold = mean_sim - 1.5 * std_sim
        except Exception as e:
            # embedding 不可用，退化到 recursive
            logger.warning("语义分块 embedding 不可用，降级为 recursive: %s", e)
            return self._chunk_text_recursive(text)

        # 按低相似度位置断块
        groups: List[List[str]] = [[sentences[0]]]
        for i, sim in enumerate(sims):
            if sim < threshold:
                groups.append([sentences[i + 1]])
            else:
                groups[-1].append(sentences[i + 1])

        # 合并成不超过 chunk_size 的块
        result: List[Dict[str, Any]] = []
        chunk_id = 0
        for group in groups:
            buf = ""
            for sent in group:
                candidate = (buf + " " + sent).strip()
                if len(candidate) <= self.chunk_size:
                    buf = candidate
                else:
                    if buf:
                        result.append({"chunk_id": f"chunk_{chunk_id}", "content": buf})
                        chunk_id += 1
                    buf = sent
            if buf:
                result.append({"chunk_id": f"chunk_{chunk_id}", "content": buf})
                chunk_id += 1

        return result if result else self._chunk_text_recursive(text)

    # ── 分层分块 ─────────────────────────────────────────────────────────────
    def _chunk_text_hierarchical(self, text: str) -> List[Dict[str, Any]]:
        """分层分块：先切父块，再对每个父块切子块；子块 metadata 含 parent_chunk_id"""
        # 父块用 recursive，使用 parent_chunk_size
        parent_chunks = self._chunk_text_recursive(
            text,
            chunk_size=self.parent_chunk_size,
            chunk_overlap=self.chunk_overlap,
        )

        all_chunks: List[Dict[str, Any]] = []
        child_counter = 0

        for parent in parent_chunks:
            parent_id = parent["chunk_id"]
            parent_content = parent["content"]

            # 对父块再切子块
            child_chunks = self._chunk_text_recursive(
                parent_content,
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
            )

            if not child_chunks:
                continue

            # 保存父块（供检索时拉取完整上下文）
            all_chunks.append({
                "chunk_id": f"parent_{parent_id}",
                "content": parent_content,
                "is_parent": True,
            })

            for child in child_chunks:
                all_chunks.append({
                    "chunk_id": f"child_{child_counter}_{child['chunk_id']}",
                    "content": child["content"],
                    "parent_chunk_id": f"parent_{parent_id}",
                    "is_parent": False,
                })
                child_counter += 1

        return all_chunks

    def _clean_text(self, text: str) -> str:
        """清理文本"""
        # 移除多余空白
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = re.sub(r" {2,}", " ", text)

        # 移除特殊字符
        text = re.sub(r"[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f]", "", text)

        return text.strip()

    def _generate_file_id(self, file_path: str, project_id: str = "") -> str:
        """生成文件 ID（使用完整路径，避免同项目同名文件碰撞）"""
        key = f"{project_id}:{Path(file_path).resolve()}"
        return hashlib.md5(key.encode()).hexdigest()[:12]

    def save_uploaded_file(
        self,
        content: bytes,
        filename: str,
        project_id: Optional[str] = None,
    ) -> str:
        """保存上传的文件"""
        # 验证文件格式
        ext = Path(filename).suffix.lower()
        if ext not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {ext}")

        # 生成唯一文件名
        file_id = hashlib.md5(content).hexdigest()[:8]
        safe_filename = f"{file_id}_{filename}"

        # 确定保存路径
        if project_id:
            project_dir = Path(self.upload_dir) / project_id
            project_dir.mkdir(parents=True, exist_ok=True)
            file_path = project_dir / safe_filename
        else:
            file_path = Path(self.upload_dir) / safe_filename

        # 保存文件
        with open(file_path, "wb") as f:
            f.write(content)

        return str(file_path)

    def is_supported(self, filename: str) -> bool:
        """检查文件是否支持"""
        ext = Path(filename).suffix.lower()
        return ext in self.supported_formats


# 全局文件处理器
file_processor = FileProcessor()
